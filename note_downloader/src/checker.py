# checker.py

import os
import json
import hashlib
# Try importing necessary libraries, handle if missing initially
try:
    from PyPDF2 import PdfReader
except ImportError:
    PdfReader = None # Mark as unavailable
try:
    from pptx import Presentation
except ImportError:
    Presentation = None # Mark as unavailable
try:
    from docx import Document
except ImportError:
    Document = None # Mark as unavailable

def file_hash(file_path):
    """对任意文件做 SHA256 哈希"""
    h = hashlib.sha256()
    try:
        with open(file_path, 'rb') as f:
            for chunk in iter(lambda: f.read(8192), b''):
                h.update(chunk)
        return h.hexdigest()
    except IOError as e:
        print(f"[ERROR] checker.file_hash:无法读取文件 {file_path}: {e}")
        return None # Return None on error

def extract_text_signature(file_path):
    """
    根据后缀提取文档文本并做 SHA256 哈希，若出错或库未安装则退回二进制哈希。
    支持 .pdf, .pptx/.ppt, .docx/.doc；其他类型直接做二进制哈希。
    """
    ext = os.path.splitext(file_path)[1].lower()
    text_to_hash = None

    try:
        if ext == ".pdf" and PdfReader:
            text = ""
            try:
                reader = PdfReader(file_path)
                for page in reader.pages:
                    page_text = page.extract_text()
                    if page_text:
                        text += page_text
                text_to_hash = text
            except Exception as e:
                print(f"[WARN] checker.extract_text_signature:处理 PDF 文件 '{os.path.basename(file_path)}' 时出错: {e}。将使用二进制哈希。")

        elif ext in [".pptx", ".ppt"] and Presentation:
            text = ""
            try:
                prs = Presentation(file_path)
                for slide in prs.slides:
                    for shape in slide.shapes:
                        if hasattr(shape, 'text'):
                            text += shape.text
                text_to_hash = text
            except Exception as e:
                print(f"[WARN] checker.extract_text_signature:处理 PPT/PPTX 文件 '{os.path.basename(file_path)}' 时出错: {e}。将使用二进制哈希。")

        elif ext in [".docx", ".doc"] and Document:
            text = ""
            try:
                doc = Document(file_path)
                for para in doc.paragraphs:
                    text += para.text
                text_to_hash = text
            except Exception as e:
                print(f"[WARN] checker.extract_text_signature:处理 DOC/DOCX 文件 '{os.path.basename(file_path)}' 时出错: {e}。将使用二进制哈希。")

        if text_to_hash is not None:
             # Ensure consistent encoding before hashing
            return hashlib.sha256(text_to_hash.encode('utf-8', errors='ignore')).hexdigest()

    except Exception as e:
         # Catch any other unexpected error during extraction attempt
        print(f"[ERROR] checker.extract_text_signature:提取文本时发生意外错误 '{os.path.basename(file_path)}': {e}。将使用二进制哈希。")


    # 如果库未安装、提取失败或非支持格式，退回二进制哈希
    # print(f"[DEBUG] checker.extract_text_signature: 对 '{os.path.basename(file_path)}' 使用二进制哈希。")
    return file_hash(file_path)


def build_manifest(root_dir, manifest_path):
    """
    遍历 root_dir 下所有文件，计算每个文件的签名，生成或更新 manifest JSON。
    返回 { 相对路径: 签名 } 字典。
    注意：键使用相对于 root_dir 的路径。
    """
    manifest = {}
    if not os.path.isdir(root_dir):
        print(f"[WARN] checker.build_manifest:根目录 '{root_dir}' 不存在。")
        return manifest

    print(f"开始构建 '{root_dir}' 的文件清单...")
    file_count = 0
    error_count = 0
    for dirpath, _, filenames in os.walk(root_dir):
        for fn in filenames:
            # 排除 manifest 文件本身
            if fn == os.path.basename(manifest_path) and os.path.abspath(dirpath) == os.path.abspath(os.path.dirname(manifest_path)):
                continue

            full_path = os.path.join(dirpath, fn)
            # 检查是否是有效文件
            if not os.path.isfile(full_path):
                continue

            try:
                rel_path = os.path.relpath(full_path, root_dir).replace('\\', '/') # Use forward slashes for consistency
                sig = extract_text_signature(full_path)
                if sig: # Only add if signature was successfully calculated
                    manifest[rel_path] = sig
                    file_count += 1
                else:
                     error_count +=1 # Increment error if signature failed
            except Exception as e:
                 print(f"[ERROR] checker.build_manifest:处理文件 '{full_path}' 时出错: {e}")
                 error_count += 1

    print(f"清单构建完成。共处理 {file_count} 个文件，{error_count} 个签名计算失败。")
    try:
        with open(manifest_path, 'w', encoding='utf-8') as f:
            json.dump(manifest, f, indent=2, ensure_ascii=False, sort_keys=True) # Sort keys for consistency
        print(f"清单已保存到: {manifest_path}")
    except IOError as e:
        print(f"[ERROR] checker.build_manifest:无法写入清单文件 {manifest_path}: {e}")

    return manifest

def load_manifest(manifest_path):
    """载入已有 manifest，若不存在或无效则返回空 dict"""
    if os.path.exists(manifest_path):
        try:
            with open(manifest_path, 'r', encoding='utf-8') as f:
                manifest = json.load(f)
            if isinstance(manifest, dict):
                 print(f"成功从 {manifest_path} 加载清单。")
                 return manifest
            else:
                 print(f"[ERROR] checker.load_manifest:清单文件 {manifest_path} 格式无效 (不是字典)。将返回空清单。")
                 return {}
        except (IOError, json.JSONDecodeError) as e:
            print(f"[ERROR] checker.load_manifest:无法加载或解析清单文件 {manifest_path}: {e}。将返回空清单。")
            return {}
    # print(f"[INFO] checker.load_manifest:清单文件 {manifest_path} 不存在。将返回空清单。")
    return {}

# Basic check for library availability
if not PdfReader:
    print("[WARN] checker.py: PyPDF2 未安装或导入失败，PDF 文本签名功能不可用。")
if not Presentation:
    print("[WARN] checker.py: python-pptx 未安装或导入失败，PPTX/PPT 文本签名功能不可用。")
if not Document:
     print("[WARN] checker.py: python-docx 未安装或导入失败，DOCX/DOC 文本签名功能不可用。")
